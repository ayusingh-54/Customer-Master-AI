"""
Customer Master Data AI - Agent Presentation Generator
Creates a professional PPTX covering all 6 AI agents
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import datetime

# ─── Colour Palette ───────────────────────────────────────────────────────────
C_DARK_BLUE   = RGBColor(0x0D, 0x2B, 0x55)   # Navy header bg
C_MID_BLUE    = RGBColor(0x1B, 0x4F, 0x9B)   # Accent / title bar
C_LIGHT_BLUE  = RGBColor(0xD6, 0xE4, 0xF7)   # Table header bg
C_ACCENT      = RGBColor(0xFF, 0x8C, 0x00)   # Orange accent (bullets, icons)
C_GREEN       = RGBColor(0x1E, 0x8B, 0x4A)   # Benefit / positive
C_RED         = RGBColor(0xC0, 0x39, 0x2B)   # Risk / negative
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_GREY   = RGBColor(0x33, 0x33, 0x33)
C_LIGHT_GREY  = RGBColor(0xF0, 0xF4, 0xF8)
C_YELLOW      = RGBColor(0xFF, 0xF0, 0xCC)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        if line_width:
            line.width = line_width
    else:
        line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, font_size=18, bold=False, color=C_DARK_GREY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=14, bold=False, color=C_DARK_GREY,
             align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False, bullet_char=""):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    if bullet_char:
        run0 = p.add_run()
        run0.text = bullet_char + "  "
        run0.font.size = Pt(font_size)
        run0.font.color.rgb = C_ACCENT
        run0.font.bold = True
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def section_header(slide, title, subtitle=""):
    """Full-width dark blue header at top"""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.3), fill_color=C_DARK_BLUE)
    add_text(slide, title, Inches(0.35), Inches(0.12), Inches(12), Inches(0.65),
             font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.35), Inches(0.75), Inches(12), Inches(0.45),
                 font_size=15, color=RGBColor(0xAD, 0xC8, 0xEE), align=PP_ALIGN.LEFT)


def add_badge(slide, label, l, t, w, h, bg=C_MID_BLUE, fg=C_WHITE, font_size=11):
    add_rect(slide, l, t, w, h, fill_color=bg)
    add_text(slide, label, l + Inches(0.05), t + Inches(0.04), w - Inches(0.1),
             h - Inches(0.08), font_size=font_size, bold=True, color=fg,
             align=PP_ALIGN.CENTER)


def two_column_box(slide, left_title, right_title, left_items, right_items,
                   l_color=C_LIGHT_BLUE, r_color=C_YELLOW,
                   lt_color=C_MID_BLUE, rt_color=RGBColor(0xB7, 0x5A, 0x00),
                   top=Inches(1.4), box_h=Inches(5.5)):
    col_w = Inches(6.2)
    gap   = Inches(0.33)
    left  = Inches(0.25)
    right = left + col_w + gap

    # Left box
    add_rect(slide, left, top, col_w, box_h, fill_color=l_color)
    add_rect(slide, left, top, col_w, Inches(0.5), fill_color=lt_color)
    add_text(slide, left_title, left + Inches(0.1), top + Inches(0.04),
             col_w - Inches(0.2), Inches(0.42), font_size=14, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    txL = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.6),
                                   col_w - Inches(0.3), box_h - Inches(0.75))
    txL.text_frame.word_wrap = True
    for item in left_items:
        add_para(txL.text_frame, item, font_size=13, bullet_char="▸")

    # Right box
    add_rect(slide, right, top, col_w, box_h, fill_color=r_color)
    add_rect(slide, right, top, col_w, Inches(0.5), fill_color=rt_color)
    add_text(slide, right_title, right + Inches(0.1), top + Inches(0.04),
             col_w - Inches(0.2), Inches(0.42), font_size=14, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    txR = slide.shapes.add_textbox(right + Inches(0.15), top + Inches(0.6),
                                   col_w - Inches(0.3), box_h - Inches(0.75))
    txR.text_frame.word_wrap = True
    for item in right_items:
        add_para(txR.text_frame, item, font_size=13, bullet_char="▸")


def three_benefit_boxes(slide, boxes, top=Inches(2.2)):
    """boxes = list of (icon, title, body) up to 3"""
    w = Inches(3.9)
    h = Inches(4.6)
    starts = [Inches(0.25), Inches(4.67), Inches(9.1)]
    colors  = [C_DARK_BLUE, C_MID_BLUE, RGBColor(0x15, 0x6E, 0x3C)]
    for i, (icon, title, body) in enumerate(boxes):
        add_rect(slide, starts[i], top, w, h, fill_color=colors[i])
        add_text(slide, icon, starts[i], top + Inches(0.15), w, Inches(0.9),
                 font_size=36, align=PP_ALIGN.CENTER, color=C_WHITE)
        add_text(slide, title, starts[i] + Inches(0.1), top + Inches(1.0),
                 w - Inches(0.2), Inches(0.5), font_size=15, bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER)
        txB = slide.shapes.add_textbox(starts[i] + Inches(0.15), top + Inches(1.6),
                                       w - Inches(0.3), h - Inches(1.8))
        txB.text_frame.word_wrap = True
        for line in body:
            add_para(txB.text_frame, line, font_size=12, color=C_WHITE, bullet_char="•")


def add_flow_boxes(slide, steps, top=Inches(2.4), h=Inches(1.2)):
    """Horizontal flow diagram"""
    n = len(steps)
    total_w = Inches(12.6)
    box_w = total_w / (n * 2 - 1)
    arrow_w = box_w
    colors_cycle = [C_MID_BLUE, RGBColor(0x1A, 0x6E, 0x9A),
                    RGBColor(0x12, 0x5E, 0x7C), C_DARK_BLUE, RGBColor(0x0A, 0x4A, 0x6A)]
    x = Inches(0.35)
    for i, step in enumerate(steps):
        c = colors_cycle[i % len(colors_cycle)]
        add_rect(slide, x, top, box_w, h, fill_color=c)
        num_box = slide.shapes.add_textbox(x + Inches(0.05), top + Inches(0.08),
                                           box_w - Inches(0.1), h - Inches(0.15))
        num_box.text_frame.word_wrap = True
        p = num_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r1 = p.add_run()
        r1.text = f"Step {i+1}\n"
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = C_ACCENT
        p2 = num_box.text_frame.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = step
        r2.font.size = Pt(11)
        r2.font.color.rgb = C_WHITE
        x += box_w
        if i < n - 1:
            # Arrow connector
            arr = slide.shapes.add_shape(1, x, top + h/4, arrow_w, h/2)
            arr.fill.solid(); arr.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            arr.line.fill.background()
            arrt = slide.shapes.add_textbox(x, top + h/4 + Inches(0.1), arrow_w, h/2)
            arrt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            r = arrt.text_frame.paragraphs[0].add_run()
            r.text = "➔"
            r.font.size = Pt(18)
            r.font.color.rgb = C_MID_BLUE
            x += arrow_w


# ─── Presentation Setup ───────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=C_DARK_BLUE)
add_rect(slide, Inches(0), Inches(2.8), SLIDE_W, Inches(0.08), fill_color=C_ACCENT)

add_text(slide, "Customer Master Data AI", Inches(0.8), Inches(0.7), Inches(11.5), Inches(1.2),
         font_size=42, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "6 Intelligent Agents — Oracle EBS TCA Automation",
         Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.7),
         font_size=22, color=RGBColor(0xAD, 0xC8, 0xEE), align=PP_ALIGN.CENTER)
add_text(slide, "How Each Agent Works • Manual vs AI • Benefits • Day-to-Day Workflow",
         Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.55),
         font_size=16, italic=True, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

# Agent badges row
agent_labels = ["1  Dedup", "2  Address", "3  Credit", "4  Contact", "5  Relations", "6  Archive"]
badge_w = Inches(1.9)
badge_x = Inches(0.45)
for i, lbl in enumerate(agent_labels):
    bg = C_ACCENT if i % 2 == 0 else C_MID_BLUE
    add_badge(slide, lbl, badge_x + i * Inches(2.05), Inches(4.1), badge_w, Inches(0.55),
              bg=bg, font_size=12)

add_text(slide, f"Prepared: {datetime.date.today().strftime('%B %Y')}   |   Powered by Claude AI + MCP",
         Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.4),
         font_size=12, italic=True, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – PROJECT OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=C_LIGHT_GREY)
section_header(slide, "Project Overview", "Customer Master Data Management — Powered by AI")

# Central architecture diagram
add_rect(slide, Inches(0.3), Inches(1.5), Inches(12.5), Inches(5.7), fill_color=C_WHITE)

# Boxes
arch_items = [
    (Inches(0.5),  Inches(1.7),  "Claude AI\n(Your Assistant)",       C_DARK_BLUE),
    (Inches(3.2),  Inches(1.7),  "MCP Server\n(server.py)",           C_MID_BLUE),
    (Inches(5.9),  Inches(1.7),  "6 Agent Modules\n(agents/*.py)",    C_GREEN),
    (Inches(8.6),  Inches(1.7),  "Oracle EBS /\nSQLite DB",           RGBColor(0x7D, 0x3C, 0x98)),
    (Inches(11.3), Inches(1.7),  "Excel &\nAudit Reports",            RGBColor(0xB7, 0x5A, 0x00)),
]
bw, bh = Inches(2.5), Inches(1.0)
for x, y, label, col in arch_items:
    add_rect(slide, x, y, bw, bh, fill_color=col)
    add_text(slide, label, x + Inches(0.1), y + Inches(0.1), bw - Inches(0.2), bh - Inches(0.2),
             font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# Arrows
for i in range(4):
    ax = Inches(0.5) + (i + 1) * Inches(2.8) - Inches(0.15)
    add_text(slide, "➔", ax, Inches(1.95), Inches(0.45), Inches(0.45),
             font_size=20, color=C_ACCENT, align=PP_ALIGN.CENTER)

# Agent grid
agent_data = [
    ("Agent 1", "Deduplication",    "Fuzzy-matches & merges duplicate customer records",   C_MID_BLUE),
    ("Agent 2", "Address Validation","Validates & geo-enriches party site addresses",       C_GREEN),
    ("Agent 3", "Credit Scoring",   "Auto-adjusts credit limits from payment history",     RGBColor(0xC0, 0x50, 0x00)),
    ("Agent 4", "Contact Mgmt",     "Manages bounced emails & missing contact points",     RGBColor(0x7D, 0x3C, 0x98)),
    ("Agent 5", "Relationships",    "Builds & maintains the corporate relationship graph",  RGBColor(0x1E, 0x60, 0x80)),
    ("Agent 6", "Archiving",        "Identifies & archives dormant customer accounts",      RGBColor(0x6E, 0x2B, 0x00)),
]
row_h = Inches(0.62)
for i, (num, name, desc, col) in enumerate(agent_data):
    row_x, row_y = (Inches(0.45) if i < 3 else Inches(6.9)), Inches(3.0) + (i % 3) * row_h
    add_rect(slide, row_x, row_y, Inches(0.6), row_h - Inches(0.06), fill_color=col)
    add_text(slide, num, row_x + Inches(0.02), row_y + Inches(0.12), Inches(0.56), Inches(0.38),
             font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, row_x + Inches(0.62), row_y, Inches(5.8), row_h - Inches(0.06), fill_color=C_WHITE)
    add_text(slide, f"{name}:  {desc}", row_x + Inches(0.72), row_y + Inches(0.1),
             Inches(5.6), row_h - Inches(0.2), font_size=12, color=C_DARK_GREY)

# ══════════════════════════════════════════════════════════════════════════════
# AGENT TEMPLATE FUNCTION
# ══════════════════════════════════════════════════════════════════════════════

def agent_slide_overview(prs, agent_num, agent_name, tagline, what_it_does,
                         tools, workflow_steps, badge_col):
    """Slide: What the agent does"""
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=C_LIGHT_GREY)
    section_header(slide, f"Agent {agent_num}:  {agent_name}", tagline)

    # Badge
    add_badge(slide, f"Agent {agent_num}", Inches(11.8), Inches(0.15), Inches(1.3), Inches(0.45),
              bg=badge_col)

    # What it does box
    add_rect(slide, Inches(0.25), Inches(1.45), Inches(7.9), Inches(5.85), fill_color=C_WHITE)
    add_rect(slide, Inches(0.25), Inches(1.45), Inches(7.9), Inches(0.5), fill_color=badge_col)
    add_text(slide, "What This Agent Does", Inches(0.35), Inches(1.48),
             Inches(7.7), Inches(0.44), font_size=14, bold=True, color=C_WHITE)

    tx = slide.shapes.add_textbox(Inches(0.4), Inches(2.05), Inches(7.6), Inches(2.5))
    tx.text_frame.word_wrap = True
    for item in what_it_does:
        add_para(tx.text_frame, item, font_size=13, bullet_char="✔")

    # Workflow steps (horizontal)
    add_text(slide, "Automated Workflow", Inches(0.35), Inches(4.0),
             Inches(7.8), Inches(0.38), font_size=13, bold=True, color=badge_col)
    step_h = Inches(0.7)
    sy = Inches(4.45)
    sw = Inches(7.6) / (len(workflow_steps) * 2 - 1)
    sx = Inches(0.35)
    for i, step in enumerate(workflow_steps):
        add_rect(slide, sx, sy, sw, step_h, fill_color=badge_col)
        stx = slide.shapes.add_textbox(sx + Inches(0.04), sy + Inches(0.05),
                                       sw - Inches(0.08), step_h - Inches(0.1))
        stx.text_frame.word_wrap = True
        p = stx.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = step; r.font.size = Pt(10); r.font.color.rgb = C_WHITE
        sx += sw
        if i < len(workflow_steps) - 1:
            arrt = slide.shapes.add_textbox(sx, sy + Inches(0.18), sw, Inches(0.35))
            arrt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            ra = arrt.text_frame.paragraphs[0].add_run()
            ra.text = "→"; ra.font.size = Pt(18); ra.font.color.rgb = badge_col
            sx += sw

    # MCP Tools panel (right)
    add_rect(slide, Inches(8.35), Inches(1.45), Inches(4.75), Inches(5.85), fill_color=C_LIGHT_GREY)
    add_rect(slide, Inches(8.35), Inches(1.45), Inches(4.75), Inches(0.5), fill_color=C_DARK_BLUE)
    add_text(slide, "MCP Tools Available", Inches(8.45), Inches(1.48),
             Inches(4.55), Inches(0.44), font_size=14, bold=True, color=C_WHITE)

    ty = Inches(2.05)
    for tool_name, tool_desc in tools:
        add_rect(slide, Inches(8.45), ty, Inches(4.55), Inches(0.62), fill_color=C_WHITE)
        add_text(slide, tool_name, Inches(8.55), ty + Inches(0.04), Inches(4.35), Inches(0.28),
                 font_size=11, bold=True, color=C_MID_BLUE)
        add_text(slide, tool_desc, Inches(8.55), ty + Inches(0.3), Inches(4.35), Inches(0.28),
                 font_size=10, color=C_DARK_GREY, italic=True)
        ty += Inches(0.7)


def agent_slide_manual_vs_ai(prs, agent_num, agent_name, manual_steps, ai_steps,
                              badge_col, time_saved, error_reduction):
    """Slide: Manual vs AI comparison"""
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=C_LIGHT_GREY)
    section_header(slide, f"Agent {agent_num}:  {agent_name} — Manual vs AI",
                   "Without the agent you do this yourself. With it, Claude does it in seconds.")
    add_badge(slide, f"Agent {agent_num}", Inches(11.8), Inches(0.15), Inches(1.3), Inches(0.45),
              bg=badge_col)

    two_column_box(
        slide,
        "❌  WITHOUT the Agent (Manual Process)",
        "✅  WITH the Agent (AI-Powered)",
        manual_steps,
        ai_steps,
        l_color=RGBColor(0xFF, 0xEC, 0xEC),
        r_color=RGBColor(0xEC, 0xFF, 0xEC),
        lt_color=C_RED,
        rt_color=C_GREEN,
        top=Inches(1.4),
        box_h=Inches(4.8),
    )

    # KPI bar at bottom
    add_rect(slide, Inches(0.25), Inches(6.35), Inches(12.85), Inches(0.95), fill_color=C_DARK_BLUE)
    kpis = [
        ("Time Saved", time_saved),
        ("Error Reduction", error_reduction),
        ("Audit Trail", "Full automatic log"),
        ("Trigger", "Ask Claude anytime"),
    ]
    kw = Inches(3.1)
    kx = Inches(0.35)
    for kname, kval in kpis:
        add_text(slide, kname, kx, Inches(6.38), kw, Inches(0.38),
                 font_size=11, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_text(slide, kval, kx, Inches(6.72), kw, Inches(0.48),
                 font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        kx += kw + Inches(0.05)


def agent_slide_benefits_daily(prs, agent_num, agent_name, benefit_boxes,
                                daily_items, badge_col):
    """Slide: Benefits + day-to-day"""
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=C_LIGHT_GREY)
    section_header(slide, f"Agent {agent_num}:  {agent_name} — Benefits & Daily Use", "")
    add_badge(slide, f"Agent {agent_num}", Inches(11.8), Inches(0.15), Inches(1.3), Inches(0.45),
              bg=badge_col)

    three_benefit_boxes(slide, benefit_boxes, top=Inches(1.45))

    # Day-to-day block
    add_rect(slide, Inches(0.25), Inches(6.2), Inches(12.85), Inches(1.1), fill_color=badge_col)
    add_text(slide, "Day-to-Day Usage", Inches(0.35), Inches(6.22),
             Inches(2.8), Inches(0.45), font_size=14, bold=True, color=C_WHITE)
    daily_str = "    |    ".join(daily_items)
    add_text(slide, daily_str, Inches(0.35), Inches(6.65), Inches(12.6), Inches(0.55),
             font_size=12, color=C_WHITE, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# AGENT 1 — DEDUPLICATION
# ══════════════════════════════════════════════════════════════════════════════
A1_COLOR = RGBColor(0x1B, 0x4F, 0x9B)

agent_slide_overview(
    prs, 1, "Deduplication Agent", "Finds & merges duplicate customer parties automatically",
    what_it_does=[
        "Scans all active ORGANIZATION and PERSON parties in the database",
        "Uses Jaro-Winkler fuzzy matching (≥88% similarity threshold) to detect near-duplicate names",
        "Also catches exact duplicates via matching tax_reference or DUNS numbers",
        "Identifies a 'Golden Record' (the master to keep) from each duplicate group",
        "Merges all linked accounts, sites, contacts, and relationships to the golden record",
        "Marks the duplicate party as 'Merged' — it is never hard-deleted for audit safety",
    ],
    tools=[
        ("find_duplicate_parties", "Fuzzy-scan all parties, return groups"),
        ("merge_duplicate_parties", "Merge duplicate into golden record"),
        ("export_duplicates_to_excel", "Generate Excel duplicate report"),
    ],
    workflow_steps=["Scan DB", "Fuzzy Match", "Score Groups", "Pick Golden", "Merge Records", "Audit Log"],
    badge_col=A1_COLOR,
)

agent_slide_manual_vs_ai(
    prs, 1, "Deduplication Agent",
    manual_steps=[
        "Open Oracle EBS, navigate to Trading Community Architecture",
        "Run a manual name search for each suspected duplicate",
        "Compare records side-by-side across multiple screens",
        "Decide which record is the 'master' — often subjective",
        "Manually reassign all open accounts, invoices, orders to master",
        "Log a change request ticket for the data team",
        "Wait for DBA to execute SQL merge scripts",
        "QA-check the merge and update documentation",
        "Repeat for hundreds of duplicates — easily 2-3 days of work",
    ],
    ai_steps=[
        "Type: 'Find all duplicate customers' in Claude",
        "Agent calls find_duplicate_parties — scans thousands of records in seconds",
        "Returns ranked groups with similarity scores and golden record suggestion",
        "Review the list and confirm merges with a single message",
        "Agent calls merge_duplicate_parties for each confirmed pair",
        "All accounts/sites/contacts re-linked automatically",
        "Export a full audit-ready Excel report with one command",
        "Total time: under 5 minutes for the entire customer base",
    ],
    badge_col=A1_COLOR,
    time_saved="~2-3 days → 5 minutes",
    error_reduction="~95% fewer missed dupes",
)

agent_slide_benefits_daily(
    prs, 1, "Deduplication Agent",
    benefit_boxes=[
        ("🔍", "Accurate Fuzzy Matching",
         ["88% Jaro-Winkler threshold", "Catches 'Acme Corp' vs 'ACME Corp.'", "Tax/DUNS exact match",
          "No missed duplicates", "Configurable threshold"]),
        ("⚡", "Instant Bulk Processing",
         ["Scans 1,000s of parties instantly", "Groups duplicates automatically", "Golden record suggested",
          "No manual comparison needed", "Runs on-demand or scheduled"]),
        ("📋", "Safe & Auditable",
         ["No hard deletes — records are 'Merged'", "Full audit_log entry per merge",
          "Excel report for sign-off", "Roll-back possible", "EBS-compliant approach"]),
    ],
    daily_items=[
        "Morning: 'Scan for new duplicates from yesterday's imports'",
        "Weekly: Export duplicate report for data governance review",
        "Post-migration: Full dedup sweep after every data load",
        "Ad-hoc: 'Is Acme Corp a duplicate of Acme Corporation?'",
    ],
    badge_col=A1_COLOR,
)

# ══════════════════════════════════════════════════════════════════════════════
# AGENT 2 — ADDRESS VALIDATION
# ══════════════════════════════════════════════════════════════════════════════
A2_COLOR = RGBColor(0x1E, 0x8B, 0x4A)

agent_slide_overview(
    prs, 2, "Address Validation & Enrichment Agent",
    "Validates address format rules and adds geo-coordinates to every party site",
    what_it_does=[
        "Retrieves all unvalidated party sites from hz_party_sites",
        "Applies country-specific format rules: US ZIP codes, state codes, UK postcodes",
        "For non-US/UK addresses, performs pass-through acceptance by default",
        "Enriches validated addresses with latitude / longitude coordinates",
        "Sets validated=1 on success or validated=0 with reason on failure",
        "Can validate a single site, all sites for one party, or all unvalidated sites in bulk",
    ],
    tools=[
        ("validate_address", "Validate by site ID, party ID, or all"),
        ("get_unvalidated_addresses", "List all pending sites"),
        ("export_customers_to_excel", "Sites sheet with validation status"),
    ],
    workflow_steps=["Fetch Sites", "Apply Rules", "Check Format", "Geo-Enrich", "Update DB", "Report"],
    badge_col=A2_COLOR,
)

agent_slide_manual_vs_ai(
    prs, 2, "Address Validation & Enrichment Agent",
    manual_steps=[
        "Export party sites to a spreadsheet",
        "Manually look up each ZIP code / postcode format rule",
        "Cross-check state codes against the official 50-state list",
        "Reject or correct invalid records one by one",
        "Use a separate geo-coding tool or service for lat/lon",
        "Copy lat/lon back into the spreadsheet row by row",
        "Re-import corrected data to EBS — risk of overwrite errors",
        "No systematic tracking of which sites were checked vs. not",
        "Repeat the entire process whenever new sites are added",
    ],
    ai_steps=[
        "Type: 'Validate all unvalidated addresses'",
        "Agent calls get_unvalidated_addresses — returns pending list instantly",
        "Calls validate_address — applies US/UK rules and geo-enriches each site",
        "Returns pass/fail per site with specific failure reasons",
        "Failed sites shown with clear error message for correction",
        "Re-run after fixing — agent picks up only the remaining invalid sites",
        "Export Excel report: Sites sheet colour-coded green (valid) / red (invalid)",
    ],
    badge_col=A2_COLOR,
    time_saved="~1 day → 2 minutes",
    error_reduction="100% consistent rules",
)

agent_slide_benefits_daily(
    prs, 2, "Address Validation & Enrichment Agent",
    benefit_boxes=[
        ("🌍", "Country-Aware Rules",
         ["US: 50-state + DC validation", "US: ZIP 5-digit / ZIP+4 regex", "UK: postcode format regex",
          "Other countries: pass-through", "Easy to extend new countries"]),
        ("📍", "Geo-Coordinate Enrichment",
         ["Adds lat/lon to every validated site", "Enables map-based analytics",
          "Powers logistics / delivery routing", "Used in Excel export map views", "Zero manual effort"]),
        ("✅", "Data Quality Assurance",
         ["Consistent pass/fail per site", "Clear failure reasons stored", "Colour-coded Excel export",
          "Validates at import time", "Prevents bad data from entering EBS"]),
    ],
    daily_items=[
        "After every customer import: 'Validate all new addresses'",
        "Weekly: Export Sites report to check outstanding invalid addresses",
        "On-boarding new customer: 'Validate the address for party 1042'",
        "Logistics review: Export customer map (lat/lon) for delivery planning",
    ],
    badge_col=A2_COLOR,
)

# ══════════════════════════════════════════════════════════════════════════════
# AGENT 3 — CREDIT SCORING
# ══════════════════════════════════════════════════════════════════════════════
A3_COLOR = RGBColor(0xC0, 0x50, 0x00)

agent_slide_overview(
    prs, 3, "Credit Limit Auto-Adjustment Agent",
    "Scores payment performance and auto-adjusts customer credit limits",
    what_it_does=[
        "Calculates a credit score for each account from 5 payment metrics",
        "Metrics: payment speed, overdue invoices, return rate, order volume, credit utilization",
        "Score ≥+2 → INCREASE limit by 15%;  Score ≤-2 → DECREASE by 25%",
        "Applies changes directly to hz_cust_accounts when approved",
        "Supports 'evaluate only' mode (no changes) for review before action",
        "run_credit_sweep() processes ALL active accounts in one operation",
    ],
    tools=[
        ("evaluate_credit", "Score + recommend (no DB change)"),
        ("apply_credit_adjustment", "Score + apply change to DB"),
        ("run_credit_sweep", "Evaluate ALL accounts at once"),
        ("export_credit_report_to_excel", "Ranked credit risk report"),
    ],
    workflow_steps=["Pull Invoices", "Score Payments", "Check Utilization", "Calculate Delta", "Apply Limit", "Log Action"],
    badge_col=A3_COLOR,
)

agent_slide_manual_vs_ai(
    prs, 3, "Credit Limit Auto-Adjustment Agent",
    manual_steps=[
        "Pull aged debt report from EBS Receivables",
        "Export orders to Excel and calculate average days-to-pay manually",
        "Count overdue invoices per customer in spreadsheet",
        "Calculate return rates using separate OE order report",
        "Compute credit utilization (balance ÷ limit) for each account",
        "Apply a subjective judgement call for each credit change",
        "Submit change requests through credit approval workflow",
        "Wait for manager sign-off — often days or weeks",
        "Update credit limits manually in EBS one account at a time",
        "Risk: inconsistent scoring, favouritism, forgotten reviews",
    ],
    ai_steps=[
        "Type: 'Run a full credit sweep for all customers'",
        "Agent runs run_credit_sweep() — evaluates every active account",
        "Returns ranked list with scores, recommendations, new limits",
        "Review exceptions: 'Show me all customers recommended for credit decrease'",
        "Approve in bulk or individually: 'Apply credit adjustment for account 1001'",
        "Export colour-coded Excel credit report for manager sign-off",
        "Entire process: under 60 seconds for the full customer base",
    ],
    badge_col=A3_COLOR,
    time_saved="~3 days → 60 seconds",
    error_reduction="100% consistent scoring",
)

agent_slide_benefits_daily(
    prs, 3, "Credit Limit Auto-Adjustment Agent",
    benefit_boxes=[
        ("📊", "Objective Scoring Model",
         ["5 weighted metrics applied equally", "No subjective human bias", "Configurable thresholds",
          "Consistent across all accounts", "Documented scoring rationale"]),
        ("💰", "Risk-Adjusted Limits",
         ["Rewards good payers with +15% increase", "Reduces exposure: -25% for bad payers",
          "Prevents over-extension of credit", "Maximises revenue from reliable customers",
          "Balances risk automatically"]),
        ("📈", "Proactive Management",
         ["Monthly sweep catches risk early", "Colour-coded Excel dashboard", "Full audit trail",
          "Can trigger alerts on score changes", "Supports credit review meetings"]),
    ],
    daily_items=[
        "Month-end: 'Run full credit sweep and export the credit report'",
        "Daily: 'Evaluate credit for the 3 new accounts added today'",
        "Collections: 'Which customers have credit score ≤ -2?'",
        "Sales: 'What is the credit utilisation for Acme Corp?'",
    ],
    badge_col=A3_COLOR,
)

# ══════════════════════════════════════════════════════════════════════════════
# AGENT 4 — CONTACT POINT MAINTENANCE
# ══════════════════════════════════════════════════════════════════════════════
A4_COLOR = RGBColor(0x6E, 0x27, 0x9A)

agent_slide_overview(
    prs, 4, "Contact Point Maintenance Agent",
    "Manages bounced emails, invalid contacts, and missing contact points",
    what_it_does=[
        "Stores and manages EMAIL, PHONE, FAX, and WEB contacts in hz_contact_points",
        "When a bounce or invalid report is received, marks the contact as status='I' (Invalid)",
        "Automatically searches for an alternate valid contact of the same type at the same party",
        "Identifies all active parties that have zero valid contact points",
        "Supports adding new contacts with duplicate prevention (unique constraint)",
        "Full contact lifecycle: Active → Invalid, or adding new contacts when needed",
    ],
    tools=[
        ("get_contact_points", "List all contacts for a party"),
        ("mark_contact_invalid", "Flag contact as invalid/bounced"),
        ("add_contact_point", "Add new EMAIL/PHONE/FAX/WEB"),
        ("get_parties_needing_contact", "Find parties with NO contacts"),
    ],
    workflow_steps=["Get Contacts", "Detect Bounce", "Mark Invalid", "Find Alternate", "Alert if None", "Log Action"],
    badge_col=A4_COLOR,
)

agent_slide_manual_vs_ai(
    prs, 4, "Contact Point Maintenance Agent",
    manual_steps=[
        "Receive bounce-back email notification from mail server",
        "Manually search EBS for which customer the email belongs to",
        "Navigate to Trading Community Architecture → Contact Points",
        "Manually deactivate the bounced email address",
        "Search for any alternate email addresses in the record",
        "If none, create a ticket for sales/account manager to find new contact",
        "No systematic visibility of which customers have zero contacts",
        "Bounced emails often sit unresolved for weeks",
        "Risk: invoices sent to invalid addresses, delayed payments",
    ],
    ai_steps=[
        "Forward bounce notification to Claude or type: 'Mark email ID 45 as bounced'",
        "Agent calls mark_contact_invalid — immediately deactivates the contact",
        "Automatically searches for alternate contact at the same party",
        "Returns result: found alternate (no action needed) or 'no contact — action required'",
        "Batch: 'Find all parties with no valid contact points'",
        "Agent returns the full list — ready for outreach campaign",
        "Add new contact: 'Add email john@acme.com for party 1001'",
        "Duplicate check runs automatically — no double entries",
    ],
    badge_col=A4_COLOR,
    time_saved="~30 min/bounce → instant",
    error_reduction="Zero unresolved bounces",
)

agent_slide_benefits_daily(
    prs, 4, "Contact Point Maintenance Agent",
    benefit_boxes=[
        ("📧", "Bounce Management",
         ["Instant invalid flagging", "Automatic alternate search", "No more stale addresses",
          "Improves invoice deliverability", "Reduces payment delays"]),
        ("📞", "Contact Coverage",
         ["Identifies zero-contact parties", "Supports 4 contact types", "Prevents duplicate contacts",
          "Enables outreach campaigns", "Always-current contact data"]),
        ("🔒", "Data Integrity",
         ["Status A/I tracked per contact", "Full audit of every change", "Unique constraint enforcement",
          "Party-level contact history", "Compliant with data governance"]),
    ],
    daily_items=[
        "On email bounce: 'Mark contact ID 45 as bounced — find alternate'",
        "Weekly: 'List all parties with no active email address'",
        "New customer: 'Add phone +44-20-1234-5678 for party 1056'",
        "Campaign prep: 'Export all parties missing contact points'",
    ],
    badge_col=A4_COLOR,
)

# ══════════════════════════════════════════════════════════════════════════════
# AGENT 5 — RELATIONSHIP GRAPH
# ══════════════════════════════════════════════════════════════════════════════
A5_COLOR = RGBColor(0x1E, 0x60, 0x80)

agent_slide_overview(
    prs, 5, "Relationship Graph Maintenance Agent",
    "Builds and maintains a directional graph of all party-to-party relationships",
    what_it_does=[
        "Maintains 5 relationship types: PARENT_SUBSIDIARY, PARTNER, RESELLER, COMPETITOR, CUSTOMER_CONTACT",
        "Each relationship is directional: subject party → object party",
        "Supports corporate restructuring: bulk-transfer all relationships to a new parent",
        "Returns the complete relationship graph for any party (all inbound + outbound links)",
        "Provides a summary view of relationship counts by type across all parties",
        "Prevents duplicate relationships via unique constraint per subject/object/type",
    ],
    tools=[
        ("get_relationships", "Graph for a single party"),
        ("add_relationship", "Create new directional link"),
        ("update_relationships_for_restructure", "Corporate restructure transfer"),
        ("get_relationship_graph_summary", "Summary by relationship type"),
    ],
    workflow_steps=["Query Party", "Fetch Graph", "Display Links", "Add/Update", "Restructure", "Summarise"],
    badge_col=A5_COLOR,
)

agent_slide_manual_vs_ai(
    prs, 5, "Relationship Graph Maintenance Agent",
    manual_steps=[
        "Manually track parent-subsidiary structure in spreadsheets",
        "Update org charts in PowerPoint whenever restructuring occurs",
        "Search EBS TCA Relationship Manager screen by screen",
        "Manually add each relationship type one at a time",
        "During M&A: manually reassign every linked record to new parent",
        "No single view of ALL relationship types in one place",
        "Spreadsheets go stale — nobody keeps them current",
        "Risk: invoicing wrong entity, wrong credit exposure view",
        "Restructure projects take days and are error-prone",
    ],
    ai_steps=[
        "Type: 'Show me all relationships for Acme Corp'",
        "Agent returns full graph: parent, subsidiaries, partners, resellers, contacts",
        "Add link: 'Create a PARENT_SUBSIDIARY relationship from TechGroup to BetaLtd'",
        "Duplicate check runs automatically — prevents double entries",
        "Restructure: 'Transfer all relationships from OldCorp to NewParent'",
        "Agent calls update_relationships_for_restructure — all links moved instantly",
        "View summary: 'How many partner and reseller relationships do we have?'",
        "All changes audited automatically",
    ],
    badge_col=A5_COLOR,
    time_saved="Restructure: days → seconds",
    error_reduction="Zero missed relationship transfers",
)

agent_slide_benefits_daily(
    prs, 5, "Relationship Graph Maintenance Agent",
    benefit_boxes=[
        ("🕸️", "Complete Graph View",
         ["5 relationship types in one place", "Directional subject→object links", "Inbound + outbound per party",
          "Parent-subsidiary hierarchies", "Partner & reseller networks"]),
        ("🔄", "Instant Restructuring",
         ["One command for M&A restructures", "All relationships transferred atomically",
          "No missed links", "Audit trail of every change", "No manual spreadsheet updates"]),
        ("📊", "Network Analytics",
         ["Graph summary by relationship type", "Identify isolated parties", "Map competitive landscape",
          "Reseller network visibility", "Informs credit & risk decisions"]),
    ],
    daily_items=[
        "Account review: 'Show me the full relationship graph for TechGroup'",
        "M&A event: 'Transfer all relationships from OldCorp to AcquirerCorp'",
        "Partnership: 'Add a RESELLER relationship between GlobalDist and BetaCo'",
        "Monthly: 'Give me a summary of all relationship types across our customer base'",
    ],
    badge_col=A5_COLOR,
)

# ══════════════════════════════════════════════════════════════════════════════
# AGENT 6 — DORMANT ACCOUNT ARCHIVING
# ══════════════════════════════════════════════════════════════════════════════
A6_COLOR = RGBColor(0x6E, 0x2B, 0x00)

agent_slide_overview(
    prs, 6, "Dormant Account Archiving Agent",
    "Identifies, classifies, and safely archives inactive customer accounts",
    what_it_does=[
        "Computes a lifecycle state for every account based on last_order_date",
        "States: PROSPECT → ACTIVE → AT-RISK → DORMANT → INACTIVE (configurable day thresholds)",
        "Finds accounts with no orders in 365+ days that are safe to archive",
        "Safety check: blocks archiving if open AR balance > $0 or recent orders exist",
        "Marks blocked accounts with clear reason codes (open AR, recent order)",
        "Supports dry-run mode: report only, no changes — ideal for review before action",
    ],
    tools=[
        ("scan_dormant_accounts", "Find archivable accounts (dry_run mode)"),
        ("sync_lifecycle_states", "Recalculate all lifecycle states"),
        ("export_dormant_report_to_excel", "Dormant accounts Excel report"),
        ("export_full_dashboard_to_excel", "Full executive dashboard"),
    ],
    workflow_steps=["Sync States", "Scan Dormant", "Check AR", "Check Orders", "Archive/Block", "Report"],
    badge_col=A6_COLOR,
)

agent_slide_manual_vs_ai(
    prs, 6, "Dormant Account Archiving Agent",
    manual_steps=[
        "Run separate reports: AR aged debt, OE order history",
        "Manually calculate days since last order for each account",
        "Cross-reference against AR balances to check outstanding debt",
        "Classify each account status manually in a spreadsheet",
        "Submit archiving requests to DBA for each dormant account",
        "No standard lifecycle state visible in EBS without custom reports",
        "Risk: archiving accounts with open invoices (write-off risk)",
        "Risk: missing AT-RISK customers who need retention outreach",
        "Process easily takes a full week for large customer bases",
    ],
    ai_steps=[
        "Type: 'Sync all lifecycle states for our customer base'",
        "Agent runs sync_lifecycle_states — every account classified instantly",
        "Type: 'Scan for dormant accounts I can safely archive (dry run)'",
        "Returns full list: safe to archive vs. blocked with clear reasons",
        "Review the report: 'Which AT-RISK customers should we call this week?'",
        "When ready: 'Archive the dormant accounts' (dry_run=False)",
        "Agent archives only safe accounts — blocked ones stay active",
        "Export full dashboard for executive review in one command",
    ],
    badge_col=A6_COLOR,
    time_saved="~1 week → 2 minutes",
    error_reduction="Zero unsafe archives",
)

agent_slide_benefits_daily(
    prs, 6, "Dormant Account Archiving Agent",
    benefit_boxes=[
        ("🔄", "Lifecycle Intelligence",
         ["5 states: PROSPECT → INACTIVE", "Configurable day thresholds", "Real-time re-calculation",
          "Drives retention outreach", "AT-RISK = revenue save opportunity"]),
        ("🛡️", "Safe Archiving",
         ["Blocks if open AR balance exists", "Blocks if recent order present", "Dry-run mode for safe preview",
          "Never archives paying customers", "Full audit of every archive action"]),
        ("📉", "DB Health & Performance",
         ["Removes inactive clutter from queries", "Reduces EBS report run time", "Clean customer lists",
          "Accurate active customer count", "Supports data retention policy"]),
    ],
    daily_items=[
        "Monday: 'Sync lifecycle states — who moved to AT-RISK this week?'",
        "Monthly: 'Dry-run dormant scan — how many accounts can we archive?'",
        "Quarter-end: 'Archive all safe dormant accounts and export dashboard'",
        "Sales: 'Show me all AT-RISK accounts for re-engagement campaign'",
    ],
    badge_col=A6_COLOR,
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE — EXCEL REPORTING SUITE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=C_LIGHT_GREY)
section_header(slide, "Excel Reporting Suite", "6 professional, colour-coded export reports — generated on-demand")

reports = [
    ("📊", "Customer Master",     "4 sheets: Accounts, Sites, Contacts, Lifecycle chart",       C_MID_BLUE),
    ("💳", "Credit Risk Report",  "All accounts ranked by score; colour-coded recommendations",  A3_COLOR),
    ("🔍", "Duplicate Report",    "Duplicate groups with AUTO-MERGE / REVIEW suggestions",       A1_COLOR),
    ("📋", "Audit Log",           "Full workflow history — filterable by agent type",             C_DARK_BLUE),
    ("💤", "Dormant Accounts",    "Archivable vs. blocked accounts with block reasons",          A6_COLOR),
    ("🏆", "Executive Dashboard", "All-in-one KPI workbook with embedded charts",               C_GREEN),
]
rw = Inches(3.85)
rh = Inches(1.3)
for i, (icon, name, desc, col) in enumerate(reports):
    rx = Inches(0.35) + (i % 3) * Inches(4.35)
    ry = Inches(1.5) + (i // 3) * Inches(1.5)
    add_rect(slide, rx, ry, rw, rh, fill_color=col)
    add_text(slide, icon + "  " + name, rx + Inches(0.1), ry + Inches(0.1),
             rw - Inches(0.2), Inches(0.5), font_size=14, bold=True, color=C_WHITE)
    add_text(slide, desc, rx + Inches(0.1), ry + Inches(0.65),
             rw - Inches(0.2), Inches(0.55), font_size=11, color=C_WHITE, italic=True)

add_text(slide, "All reports: timestamped filenames, professional styling, frozen panes, alternating rows, embedded charts",
         Inches(0.35), Inches(4.75), Inches(12.6), Inches(0.5),
         font_size=13, italic=True, color=C_DARK_GREY, align=PP_ALIGN.CENTER)

add_text(slide, "Command:  export_full_dashboard_to_excel  →  One file with everything",
         Inches(0.35), Inches(5.25), Inches(12.6), Inches(0.6),
         font_size=15, bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE — DAY-IN-THE-LIFE WORKFLOW
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=C_LIGHT_GREY)
section_header(slide, "A Day in the Life — Working with All 6 Agents", "How you interact with the AI in your daily Oracle EBS data management routine")

time_blocks = [
    ("🌅  MORNING\n8:00 AM",
     ["'Sync all lifecycle states'",
      "'Scan for new duplicates from last night's import'",
      "'Validate all new customer addresses'",
      "Review AI summary — no manual report-pulling"],
     C_DARK_BLUE),
    ("☀️  MID-DAY\n12:00 PM",
     ["'Mark email for contact 45 as bounced'",
      "'Find alternate contact for Acme Corp'",
      "'Add new RESELLER relationship for BetaCo'",
      "'Evaluate credit for account 1001'"],
     C_MID_BLUE),
    ("🌆  AFTERNOON\n3:00 PM",
     ["'Run credit sweep and export report'",
      "'Which customers are AT-RISK this week?'",
      "'Dry-run dormant archive scan'",
      "Share Excel dashboard with management"],
     A3_COLOR),
    ("📋  MONTH-END",
     ["'Archive all safe dormant accounts'",
      "'Apply all recommended credit adjustments'",
      "'Export full executive dashboard'",
      "Data governance sign-off ready"],
     C_GREEN),
]

tw = Inches(3.0)
for i, (time_label, actions, col) in enumerate(time_blocks):
    tx = Inches(0.25) + i * Inches(3.27)
    add_rect(slide, tx, Inches(1.4), tw, Inches(5.8), fill_color=col)
    add_text(slide, time_label, tx + Inches(0.1), Inches(1.45),
             tw - Inches(0.2), Inches(0.8), font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    block_tf = slide.shapes.add_textbox(tx + Inches(0.1), Inches(2.35), tw - Inches(0.2), Inches(4.6))
    block_tf.text_frame.word_wrap = True
    for action in actions:
        add_para(block_tf.text_frame, action, font_size=12, color=C_WHITE, bullet_char="→", space_before=Pt(8))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE — SUMMARY & CALL TO ACTION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=C_DARK_BLUE)
add_rect(slide, Inches(0), Inches(2.6), SLIDE_W, Inches(0.08), fill_color=C_ACCENT)

add_text(slide, "Customer Master Data AI", Inches(0.6), Inches(0.3),
         Inches(12), Inches(0.9), font_size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "6 Agents — One Conversation — Zero Manual Data Work",
         Inches(0.6), Inches(1.15), Inches(12), Inches(0.6),
         font_size=20, italic=True, color=RGBColor(0xAD, 0xC8, 0xEE), align=PP_ALIGN.CENTER)

summary_items = [
    ("Agent 1", "Deduplication",  "Days of manual merging → 5 minutes",    A1_COLOR),
    ("Agent 2", "Address Validation", "Days of format checking → 2 minutes", A2_COLOR),
    ("Agent 3", "Credit Scoring", "3-day review cycle → 60 seconds",        A3_COLOR),
    ("Agent 4", "Contact Mgmt",   "30 min/bounce → instant resolution",      A4_COLOR),
    ("Agent 5", "Relationships",  "M&A restructure: days → seconds",         A5_COLOR),
    ("Agent 6", "Archiving",      "1-week lifecycle review → 2 minutes",     A6_COLOR),
]
sw = Inches(1.9)
sx = Inches(0.45)
for i, (num, name, saving, col) in enumerate(summary_items):
    sy = Inches(2.9) + (i % 3) * Inches(1.1) if i < 3 else Inches(2.9) + (i - 3) * Inches(1.1)
    if i == 3:
        sx = Inches(7.0)
    add_rect(slide, sx, sy, sw * 1.8, Inches(1.0), fill_color=col)
    add_text(slide, f"{num}: {name}", sx + Inches(0.1), sy + Inches(0.08),
             sw * 1.8 - Inches(0.2), Inches(0.4), font_size=12, bold=True, color=C_WHITE)
    add_text(slide, saving, sx + Inches(0.1), sy + Inches(0.52),
             sw * 1.8 - Inches(0.2), Inches(0.38), font_size=10, color=C_ACCENT, italic=True)

add_text(slide, "🚀  Start with Claude Code + MCP Server  →  Ask your first question today",
         Inches(0.6), Inches(6.5), Inches(12), Inches(0.7),
         font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ─── Save ─────────────────────────────────────────────────────────────────────
out_path = r"C:\Users\ayusi\Desktop\Customer_Master_AI_Agents_Presentation.pptx"
prs.save(out_path)
print(f"✅  Presentation saved to: {out_path}")
print(f"    Slides: {len(prs.slides)}")
